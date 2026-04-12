#!/usr/bin/env python3
"""
Aurora v2 UI Essentialist Design — PPT Mockup Generator
Generates a visual presentation of the single-screen proof system UI.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# ── Color palette ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)        # dark navy background
BG_PANEL = RGBColor(0x16, 0x21, 0x3E)        # panel background
BG_CARD = RGBColor(0x1F, 0x2B, 0x4D)         # card background
BG_INPUT = RGBColor(0x0F, 0x17, 0x2A)        # input field background
ACCENT_BLUE = RGBColor(0x4E, 0x9A, 0xF5)     # primary accent
ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0xB0)    # corroborated
ACCENT_RED = RGBColor(0xE0, 0x6C, 0x75)      # contradiction/weakened
ACCENT_YELLOW = RGBColor(0xE5, 0xC0, 0x7B)   # warning/in-progress
ACCENT_PURPLE = RGBColor(0xC6, 0x78, 0xDD)   # sincerity axis
TEXT_WHITE = RGBColor(0xE0, 0xE0, 0xE0)
TEXT_DIM = RGBColor(0x8B, 0x8B, 0xA0)
TEXT_BRIGHT = RGBColor(0xFF, 0xFF, 0xFF)
LAYER_COLORS = [
    RGBColor(0x4E, 0x9A, 0xF5),  # L1 blue
    RGBColor(0x4E, 0xC9, 0xB0),  # L2 teal
    RGBColor(0xE5, 0xC0, 0x7B),  # L3 yellow
    RGBColor(0xE0, 0x6C, 0x75),  # L4 red
    RGBColor(0xC6, 0x78, 0xDD),  # L5 purple
    RGBColor(0xFF, 0x92, 0x4C),  # L6 orange
    RGBColor(0x56, 0xB6, 0xC2),  # L7 cyan
]

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Minimal rounding
    shape.adjustments[0] = 0.02
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Pretendard"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, default_size=11,
                       default_color=TEXT_WHITE, font_name="Pretendard"):
    """lines: list of (text, size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(2)
    return txBox


# ════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)

    # Decorative line
    add_rect(slide, Inches(2), Inches(2.8), Inches(9.333), Pt(2), ACCENT_BLUE)

    add_text_box(slide, Inches(2), Inches(1.2), Inches(9.333), Inches(1),
                 "Aurora v2", 44, ACCENT_BLUE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(1.9), Inches(9.333), Inches(0.8),
                 "Self-Proving System UI Design", 24, TEXT_DIM, False, PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(3.2), Inches(9.333), Inches(0.8),
                 "Essentialist Single-Screen Proof Architecture", 20, TEXT_WHITE, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(3.8), Inches(9.333), Inches(0.6),
                 "본질주의 단일 화면 증명 체계", 16, TEXT_DIM, False, PP_ALIGN.CENTER)

    # Bottom info
    add_multiline_text(slide, Inches(2), Inches(5.5), Inches(9.333), Inches(1.5), [
        ("2016 DIDC North Korean Hacking Incident", 14, TEXT_DIM, False, PP_ALIGN.CENTER),
        ("Organized Cover-up by Defense Organizations — 7-Layer Proof System", 14, TEXT_DIM, False, PP_ALIGN.CENTER),
        ("", 8, TEXT_DIM),
        ("2026-04-11  |  Claude Opus 4.6", 11, TEXT_DIM, False, PP_ALIGN.CENTER),
    ])


# ════════════════════════════════════════
# SLIDE 2 — v1 vs v2 Paradigm Shift
# ════════════════════════════════════════
def slide_paradigm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
                 "Design Philosophy: v1 vs v2", 28, ACCENT_BLUE, True)
    add_rect(slide, Inches(0.6), Inches(0.85), Inches(12), Pt(1), ACCENT_BLUE)

    # v1 column
    add_rect(slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(5.8), BG_PANEL,
             ACCENT_RED, Pt(1))
    add_text_box(slide, Inches(0.9), Inches(1.4), Inches(5.2), Inches(0.5),
                 "Aurora v1 — Database Browser", 18, ACCENT_RED, True)

    v1_items = [
        ("First Screen", "Empty search bar + graph node cluster", "User: 'What am I looking at?'"),
        ("Default Mode", "Graph exploration", "Nodes without context = noise"),
        ("Evidence", "Separate screen", "Must leave current view to verify"),
        ("6 Modes", "Ask/Trace/Contradict/Timeline/Layer/Falsify", "Cognitive overload"),
        ("3-Axis Truth", "Existed in code only", "User never saw it"),
        ("Role-Anchors", "Not reflected", "WHO mattered, not WHICH POSITION"),
    ]
    y = Inches(2.1)
    for title, desc, problem in v1_items:
        add_multiline_text(slide, Inches(0.9), y, Inches(5.2), Inches(0.7), [
            (f"  {title}", 11, TEXT_WHITE, True),
            (f"    {desc}", 9, TEXT_DIM),
            (f"    {problem}", 9, ACCENT_RED),
        ])
        y += Inches(0.85)

    # v2 column
    add_rect(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.8), BG_PANEL,
             ACCENT_GREEN, Pt(1))
    add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5.2), Inches(0.5),
                 "Aurora v2 — Self-Proving System", 18, ACCENT_GREEN, True)

    v2_items = [
        ("First Screen", "7-Layer proof declaration + progress", "System speaks first"),
        ("Default Mode", "Contradiction pairs", "Proof, not exploration"),
        ("Evidence", "Always on the right panel", "Zero distance to proof"),
        ("No Modes", "Contextual transitions via clicks", "Natural flow"),
        ("3-Axis Truth", "Visual badges on every atom", "Truthfulness/Validity/Sincerity visible"),
        ("Role-Anchors", "Graph node labels show (role-N)", "WHICH POSITION reveals the crime"),
    ]
    y = Inches(2.1)
    for title, desc, benefit in v2_items:
        add_multiline_text(slide, Inches(7.2), y, Inches(5.2), Inches(0.7), [
            (f"  {title}", 11, TEXT_WHITE, True),
            (f"    {desc}", 9, TEXT_DIM),
            (f"    {benefit}", 9, ACCENT_GREEN),
        ])
        y += Inches(0.85)


# ════════════════════════════════════════
# SLIDE 3 — Three Essentialist Principles
# ════════════════════════════════════════
def slide_principles(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
                 "3 Essentialist Principles", 28, ACCENT_BLUE, True)
    add_text_box(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.4),
                 "본질주의 설계 3원칙", 14, TEXT_DIM)
    add_rect(slide, Inches(0.6), Inches(1.15), Inches(12), Pt(1), ACCENT_BLUE)

    principles = [
        ("01", "The System Speaks First",
         "시스템이 먼저 말한다",
         "The landing screen declares:\n'This is what is being proven.'\nNot an empty search bar — a proof statement.",
         ACCENT_BLUE),
        ("02", "Contradiction is Default",
         "모순이 기본이다",
         "The minimum unit of proof delivery is\nnot a single claim — it is an opposing pair.\nv1's default was a graph. v2's default is a contradiction.",
         ACCENT_RED),
        ("03", "Zero Distance to Evidence",
         "증거와의 거리는 0",
         "The scanned original (Record No.) is always\nvisible beside the claim. No screen transitions.\nSelf-proving = instantly verifiable.",
         ACCENT_GREEN),
    ]

    for i, (num, title_en, title_kr, desc, color) in enumerate(principles):
        x = Inches(0.6 + i * 4.1)
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.4), Inches(1.6), Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = TEXT_BRIGHT
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = False

        add_text_box(slide, x, Inches(2.6), Inches(3.8), Inches(0.5),
                     title_en, 18, TEXT_BRIGHT, True, PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(3.05), Inches(3.8), Inches(0.4),
                     title_kr, 13, TEXT_DIM, False, PP_ALIGN.CENTER)

        # Description card
        card = add_rect(slide, x + Inches(0.2), Inches(3.6), Inches(3.4), Inches(2.8), BG_CARD, color, Pt(1))
        add_multiline_text(slide, x + Inches(0.5), Inches(3.8), Inches(2.9), Inches(2.4),
                           [(line, 12, TEXT_WHITE) for line in desc.split("\n")])


# ════════════════════════════════════════
# SLIDE 4 — Landing Screen Mockup
# ════════════════════════════════════════
def slide_landing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.5),
                 "Landing Screen — First Encounter (3 seconds)", 24, ACCENT_BLUE, True)
    add_rect(slide, Inches(0.6), Inches(0.6), Inches(12), Pt(1), ACCENT_BLUE)

    # Main frame
    add_rect(slide, Inches(0.8), Inches(0.9), Inches(11.7), Inches(6.2), BG_PANEL,
             RGBColor(0x30, 0x40, 0x60), Pt(1))

    # Title area
    add_text_box(slide, Inches(1.2), Inches(1.0), Inches(10.9), Inches(0.6),
                 "2016 DIDC North Korean Hacking Incident", 22, TEXT_BRIGHT, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.2), Inches(1.5), Inches(10.9), Inches(0.4),
                 "Organized Cover-up by Defense Organizations  |  7-Layer Proof System", 13, TEXT_DIM, False, PP_ALIGN.CENTER)

    # 7-Layer tower
    layer_names_kr = [
        "Active-X 제거 사업 간 舊KIATIS 이력 제거",
        "新KIATIS 사업 추진체계 및 장교 자력 조작",
        "국전원 전속 후 SW사업관리 착수·종결",
        "新KIATIS 개발·운영·시험평가 조작",
        "허위 갑질 신고와 조작 감사",
        "군 검찰단의 사기 수사와 범죄자 낙인",
        "진정서·수사 촉구 후 기소유예",
    ]
    layer_subtitles = [
        "(DIDC 해킹 근원서버 은폐의 출발점)",
        "",
        "(국방정보화카르텔 공모 구조)",
        "(표적수사 설계의 기반)",
        "(한지훈 중령 인권침해·고립화)",
        "(증거 인멸·문서 조작)",
        "(국방부의 지속적 범죄 정당화)",
    ]
    progress = [0.90, 0.70, 0.45, 0.45, 0.85, 0.75, 0.40]  # visual progress

    tower_left = Inches(2.0)
    tower_width = Inches(9.3)
    bar_height = Inches(0.55)

    for i in range(7):
        layer_num = i + 1
        y = Inches(6.15) - i * Inches(0.63)
        color = LAYER_COLORS[i]

        # Layer number box
        num_box = add_rect(slide, tower_left, y, Inches(0.45), bar_height, color)
        tf = num_box.text_frame
        tf.paragraphs[0].text = str(layer_num)
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = TEXT_BRIGHT
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Layer name
        add_text_box(slide, tower_left + Inches(0.55), y + Inches(0.02),
                     Inches(5.5), bar_height,
                     layer_names_kr[i], 11, TEXT_WHITE, True)

        if layer_subtitles[i]:
            add_text_box(slide, tower_left + Inches(5.8), y + Inches(0.05),
                         Inches(2.2), bar_height,
                         layer_subtitles[i], 8, TEXT_DIM)

        # Progress bar background
        bar_x = tower_left + Inches(8.0)
        bar_w = Inches(1.1)
        add_rect(slide, bar_x, y + Inches(0.15), bar_w, Inches(0.25),
                 RGBColor(0x2A, 0x2A, 0x3A))
        # Progress bar fill
        fill_w = int(bar_w * progress[i])
        if fill_w > 0:
            add_rect(slide, bar_x, y + Inches(0.15), fill_w, Inches(0.25), color)

    # Stats line
    add_text_box(slide, Inches(2.0), Inches(6.5), Inches(9.3), Inches(0.3),
                 "112 atoms   |   13,495 evidence records   |   42 pseudonymized persons   |   1 open contradiction",
                 11, TEXT_DIM, False, PP_ALIGN.CENTER)

    # Search input
    add_rect(slide, Inches(3.0), Inches(6.85), Inches(7.3), Inches(0.45), BG_INPUT,
             RGBColor(0x40, 0x50, 0x70), Pt(1))
    add_text_box(slide, Inches(3.3), Inches(6.88), Inches(6.7), Inches(0.4),
                 "Ask anything, or click a Layer above.", 12, TEXT_DIM)


# ════════════════════════════════════════
# SLIDE 5 — 3-Panel Proof Screen
# ════════════════════════════════════════
def slide_proof_screen(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.45),
                 "3-Panel Proof Screen — Layer Selected", 24, ACCENT_BLUE, True)
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(12), Pt(1), ACCENT_BLUE)

    # ── LEFT PANEL: Layer Navigator ──
    lp_x, lp_w = Inches(0.3), Inches(1.6)
    add_rect(slide, lp_x, Inches(0.8), lp_w, Inches(6.3), BG_PANEL,
             RGBColor(0x30, 0x40, 0x60), Pt(1))
    add_text_box(slide, lp_x + Inches(0.1), Inches(0.9), lp_w - Inches(0.2), Inches(0.3),
                 "LAYERS", 10, TEXT_DIM, True, PP_ALIGN.CENTER)

    for i in range(7):
        y = Inches(5.9) - i * Inches(0.65)
        color = LAYER_COLORS[i]
        is_selected = (i == 3)  # Layer 4 selected

        if is_selected:
            box = add_rect(slide, lp_x + Inches(0.1), y, Inches(1.4), Inches(0.55),
                           color, color, Pt(2))
        else:
            box = add_rect(slide, lp_x + Inches(0.1), y, Inches(1.4), Inches(0.55),
                           BG_CARD)

        tf = box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = f"L{i+1}"
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = TEXT_BRIGHT if is_selected else TEXT_DIM
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Query input at bottom
    add_rect(slide, lp_x + Inches(0.1), Inches(6.5), Inches(1.4), Inches(0.45), BG_INPUT,
             RGBColor(0x40, 0x50, 0x70), Pt(1))
    add_text_box(slide, lp_x + Inches(0.15), Inches(6.52), Inches(1.3), Inches(0.4),
                 "Ask...", 9, TEXT_DIM)

    # ── CENTER PANEL: Proof Body ──
    cp_x, cp_w = Inches(2.1), Inches(6.8)
    add_rect(slide, cp_x, Inches(0.8), cp_w, Inches(6.3), BG_PANEL,
             RGBColor(0x30, 0x40, 0x60), Pt(1))

    add_text_box(slide, cp_x + Inches(0.2), Inches(0.9), cp_w - Inches(0.4), Inches(0.4),
                 "Layer 4 — 新KIATIS 개발·운영·시험평가 전·중·후 조작", 16, LAYER_COLORS[3], True)
    add_text_box(slide, cp_x + Inches(0.2), Inches(1.25), cp_w - Inches(0.4), Inches(0.3),
                 "표적수사 설계의 기반  |  28% coverage  |  14 atoms", 10, TEXT_DIM)

    # Contradiction card 1
    card1_y = Inches(1.7)
    add_rect(slide, cp_x + Inches(0.2), card1_y, cp_w - Inches(0.4), Inches(2.5),
             BG_CARD, ACCENT_RED, Pt(1))

    add_text_box(slide, cp_x + Inches(0.4), card1_y + Inches(0.1), Inches(4.5), Inches(0.3),
                 "C-L4-01  Contradiction", 12, ACCENT_RED, True)
    add_text_box(slide, cp_x + Inches(5.0), card1_y + Inches(0.1), Inches(1.5), Inches(0.3),
                 "CORROBORATED", 10, ACCENT_GREEN, True, PP_ALIGN.RIGHT)

    # Claim A
    add_rect(slide, cp_x + Inches(0.4), card1_y + Inches(0.5), Inches(2.9), Inches(1.2),
             RGBColor(0x15, 0x25, 0x45), ACCENT_BLUE, Pt(1))
    add_multiline_text(slide, cp_x + Inches(0.55), card1_y + Inches(0.55), Inches(2.6), Inches(1.0), [
        ("Claim A (Wiki)", 10, ACCENT_BLUE, True),
        ("", 4, TEXT_DIM),
        ("제57조 OT&E 환경 조항이", 10, TEXT_WHITE),
        ("제2398호에서 삽입, 제2436호에서", 10, TEXT_WHITE),
        ("삭제, 제2842호에서 재삽입됨", 10, TEXT_WHITE),
        ("(3-step flip-flop)", 9, TEXT_DIM),
    ])

    # VS
    add_text_box(slide, cp_x + Inches(3.3), card1_y + Inches(0.9), Inches(0.5), Inches(0.3),
                 "vs", 14, ACCENT_RED, True, PP_ALIGN.CENTER)

    # Claim B
    add_rect(slide, cp_x + Inches(3.8), card1_y + Inches(0.5), Inches(2.6), Inches(1.2),
             RGBColor(0x15, 0x25, 0x45), ACCENT_YELLOW, Pt(1))
    add_multiline_text(slide, cp_x + Inches(3.95), card1_y + Inches(0.55), Inches(2.3), Inches(1.0), [
        ("Claim B (Book)", 10, ACCENT_YELLOW, True),
        ("", 4, TEXT_DIM),
        ("책 비교표는 제2398호의", 10, TEXT_WHITE),
        ("Article 57을 제2129호와", 10, TEXT_WHITE),
        ("동일하다고 표기", 10, TEXT_WHITE),
    ])

    # Resolution
    add_multiline_text(slide, cp_x + Inches(0.55), card1_y + Inches(1.85), Inches(5.8), Inches(0.6), [
        ("Resolution: blind 재측정으로 wiki 확인 — 책 비교표의 column 누락 확인", 10, ACCENT_GREEN),
    ])

    # 3-axis badges
    axis_y = Inches(4.4)
    add_text_box(slide, cp_x + Inches(0.4), axis_y, Inches(2.0), Inches(0.3),
                 "진리성 (Truthfulness)  ●●●●○", 11, ACCENT_BLUE, True)
    add_text_box(slide, cp_x + Inches(2.5), axis_y, Inches(2.0), Inches(0.3),
                 "타당성 (Validity)  ●●●○○", 11, ACCENT_GREEN, True)
    add_text_box(slide, cp_x + Inches(4.6), axis_y, Inches(2.0), Inches(0.3),
                 "진실성 (Sincerity)  ●●●●●", 11, ACCENT_PURPLE, True)

    # Contradiction card 2 (preview)
    card2_y = Inches(4.9)
    add_rect(slide, cp_x + Inches(0.2), card2_y, cp_w - Inches(0.4), Inches(1.0),
             BG_CARD, RGBColor(0x40, 0x50, 0x70), Pt(1))
    add_text_box(slide, cp_x + Inches(0.4), card2_y + Inches(0.1), Inches(5.5), Inches(0.3),
                 "C-L4-02  시험평가 분리 원칙 vs 통합 원칙 뒤집기", 12, TEXT_DIM, True)
    add_text_box(slide, cp_x + Inches(0.4), card2_y + Inches(0.4), Inches(5.5), Inches(0.3),
                 "제2129호 분리 → 제2436호 통합 → 원칙 자체의 180도 전환", 10, TEXT_DIM)
    add_text_box(slide, cp_x + Inches(5.0), card2_y + Inches(0.1), Inches(1.5), Inches(0.3),
                 "NEEDS_MORE", 10, ACCENT_YELLOW, True, PP_ALIGN.RIGHT)

    # Scroll hint
    add_text_box(slide, cp_x + Inches(0.2), Inches(6.15), cp_w - Inches(0.4), Inches(0.3),
                 "  +4 more contradictions in Layer 4", 10, TEXT_DIM, False, PP_ALIGN.CENTER)

    # ── RIGHT PANEL: Evidence Context ──
    rp_x, rp_w = Inches(9.1), Inches(3.9)
    add_rect(slide, rp_x, Inches(0.8), rp_w, Inches(6.3), BG_PANEL,
             RGBColor(0x30, 0x40, 0x60), Pt(1))

    add_text_box(slide, rp_x + Inches(0.2), Inches(0.9), rp_w - Inches(0.4), Inches(0.3),
                 "EVIDENCE CONTEXT", 10, TEXT_DIM, True)

    # Record No. preview
    add_rect(slide, rp_x + Inches(0.2), Inches(1.3), rp_w - Inches(0.4), Inches(2.5),
             BG_CARD, ACCENT_BLUE, Pt(1))
    add_text_box(slide, rp_x + Inches(0.4), Inches(1.4), rp_w - Inches(0.8), Inches(0.3),
                 "Record No. 10,347", 14, ACCENT_BLUE, True)
    add_text_box(slide, rp_x + Inches(0.4), Inches(1.7), rp_w - Inches(0.8), Inches(0.3),
                 "Layer 4  |  2022-09-20  |  L4-001", 9, TEXT_DIM)

    # Scan preview placeholder
    add_rect(slide, rp_x + Inches(0.4), Inches(2.1), rp_w - Inches(0.8), Inches(1.4),
             RGBColor(0x25, 0x30, 0x50))
    add_text_box(slide, rp_x + Inches(0.4), Inches(2.5), rp_w - Inches(0.8), Inches(0.5),
                 "[Scanned PDF Preview]\n\"또는 이와 유사한 환경에서\"", 11, TEXT_DIM, False, PP_ALIGN.CENTER)

    # Related atoms
    add_text_box(slide, rp_x + Inches(0.2), Inches(4.0), rp_w - Inches(0.4), Inches(0.3),
                 "Citing Atoms", 11, TEXT_WHITE, True)
    related = [
        "2398-2842ho-otne-hedge-flipflop",
        "article-58-separation-to-integration",
        "prosecution-misapplies-2129-art-58",
    ]
    for j, atom in enumerate(related):
        add_text_box(slide, rp_x + Inches(0.3), Inches(4.35) + j * Inches(0.3),
                     rp_w - Inches(0.5), Inches(0.25),
                     f"  {atom}", 9, ACCENT_BLUE)

    # Related persons
    add_text_box(slide, rp_x + Inches(0.2), Inches(5.35), rp_w - Inches(0.4), Inches(0.3),
                 "Related Persons", 11, TEXT_WHITE, True)
    persons = [
        ("박성호", "2016해킹당시원장-1"),
        ("이지영", "공문결재자-1"),
        ("장우진", "사업실무자-1"),
    ]
    for j, (name, role) in enumerate(persons):
        add_text_box(slide, rp_x + Inches(0.3), Inches(5.7) + j * Inches(0.3),
                     rp_w - Inches(0.5), Inches(0.25),
                     f"  {name} ({role})", 10, TEXT_WHITE)

    # Graph button
    add_rect(slide, rp_x + Inches(0.5), Inches(6.7), rp_w - Inches(1.0), Inches(0.35),
             ACCENT_BLUE)
    add_text_box(slide, rp_x + Inches(0.5), Inches(6.72), rp_w - Inches(1.0), Inches(0.3),
                 "View Relationship Graph", 11, TEXT_BRIGHT, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════
# SLIDE 6 — Graph Modal (On-Demand)
# ════════════════════════════════════════
def slide_graph_modal(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.45),
                 "Graph View — On-Demand Modal Overlay", 24, ACCENT_BLUE, True)
    add_text_box(slide, Inches(0.6), Inches(0.55), Inches(8), Inches(0.3),
                 "Graph appears only when needed — not the default view", 12, TEXT_DIM)
    add_rect(slide, Inches(0.6), Inches(0.85), Inches(12), Pt(1), ACCENT_BLUE)

    # Dimmed background (simulating overlay)
    add_rect(slide, Inches(0.3), Inches(1.1), Inches(12.7), Inches(6.1),
             RGBColor(0x10, 0x10, 0x20))

    # Modal frame
    modal_x, modal_y = Inches(1.5), Inches(1.3)
    modal_w, modal_h = Inches(10.3), Inches(5.7)
    add_rect(slide, modal_x, modal_y, modal_w, modal_h, BG_PANEL,
             ACCENT_BLUE, Pt(2))

    add_text_box(slide, modal_x + Inches(0.3), modal_y + Inches(0.2),
                 Inches(6), Inches(0.4),
                 "Relationship Network — Layer 4", 18, TEXT_BRIGHT, True)
    # Close button
    add_text_box(slide, modal_x + modal_w - Inches(0.8), modal_y + Inches(0.2),
                 Inches(0.5), Inches(0.4),
                 "X", 16, TEXT_DIM, True, PP_ALIGN.RIGHT)

    # Graph area
    graph_x = modal_x + Inches(0.3)
    graph_y = modal_y + Inches(0.8)
    graph_w = Inches(6.5)
    graph_h = Inches(4.5)
    add_rect(slide, graph_x, graph_y, graph_w, graph_h, BG_DARK)

    # Nodes
    nodes = [
        (Inches(1.5), Inches(1.0), "박성호\n(원장-1)", LAYER_COLORS[0]),
        (Inches(4.5), Inches(0.5), "이태호\n(위원장-1)", LAYER_COLORS[3]),
        (Inches(1.5), Inches(3.0), "이지영\n(결재자-1)", LAYER_COLORS[2]),
        (Inches(4.5), Inches(3.0), "이준호\n(공모자-1)", LAYER_COLORS[5]),
        (Inches(3.0), Inches(1.8), "KIATIS\n성능개선사업", ACCENT_YELLOW),
    ]

    for dx, dy, label, color in nodes:
        nx, ny = graph_x + dx, graph_y + dy
        node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      nx, ny, Inches(1.2), Inches(0.7))
        node.fill.solid()
        node.fill.fore_color.rgb = color
        node.line.fill.background()
        node.adjustments[0] = 0.15
        tf = node.text_frame
        tf.word_wrap = True
        for j, line in enumerate(label.split("\n")):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(9)
            p.font.color.rgb = TEXT_BRIGHT
            p.font.bold = (j == 0)
            p.alignment = PP_ALIGN.CENTER

    # Edge labels (simplified as text)
    edges = [
        (Inches(2.8), Inches(0.6), "시험평가 지시", 8),
        (Inches(1.5), Inches(2.2), "지시", 8),
        (Inches(3.0), Inches(3.2), "공모", 8),
        (Inches(2.5), Inches(1.1), "관여", 8),
    ]
    for dx, dy, label, size in edges:
        add_text_box(slide, graph_x + dx, graph_y + dy, Inches(1.2), Inches(0.25),
                     f"--- {label} -->", size, TEXT_DIM, False, PP_ALIGN.CENTER)

    # Legend
    legend_x = modal_x + Inches(7.2)
    legend_y = modal_y + Inches(0.8)
    add_rect(slide, legend_x, legend_y, Inches(2.7), Inches(4.5), BG_CARD)

    add_text_box(slide, legend_x + Inches(0.2), legend_y + Inches(0.1),
                 Inches(2.3), Inches(0.3), "Legend", 12, TEXT_WHITE, True)

    legend_items = [
        ("Person (Role-Anchor)", TEXT_WHITE),
        ("  Node label = pseudonym", TEXT_DIM),
        ("  Sub-label = (role-N)", TEXT_DIM),
        ("", TEXT_DIM),
        ("Interaction", TEXT_WHITE),
        ("  Click node = show atoms", ACCENT_BLUE),
        ("  Double-click = navigate", ACCENT_BLUE),
        ("", TEXT_DIM),
        ("Color = Layer origin", TEXT_WHITE),
        ("  Cross-layer links visible", TEXT_DIM),
        ("  Multi-layer actors highlighted", ACCENT_RED),
    ]
    for j, (text, color) in enumerate(legend_items):
        add_text_box(slide, legend_x + Inches(0.2), legend_y + Inches(0.5) + j * Inches(0.32),
                     Inches(2.3), Inches(0.3), text, 9, color)

    # Bottom note
    add_text_box(slide, modal_x + Inches(0.3), modal_y + Inches(5.35),
                 modal_w - Inches(0.6), Inches(0.3),
                 "Powered by Cytoscape.js  |  Reused from Aurora v1 graph.component.ts (1,464 LoC)  |  Node click returns to Proof Screen",
                 9, TEXT_DIM, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════
# SLIDE 7 — User Journey Flow
# ════════════════════════════════════════
def slide_user_journey(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.45),
                 "User Journey — From First Encounter to Deep Proof", 24, ACCENT_BLUE, True)
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(12), Pt(1), ACCENT_BLUE)

    steps = [
        ("3s", "Landing", "7-Layer tower\n+ progress bars\n+ stats",
         "System declares:\n'This is what is\nbeing proven'", ACCENT_BLUE),
        ("30s", "Layer Click", "Contradiction\npairs unfold\nin center panel",
         "User reads\nopposing claims\nside by side", LAYER_COLORS[3]),
        ("10s", "Record Click", "Scanned PDF\nappears in\nright panel",
         "User verifies:\n'This is real'\n(zero distance)", ACCENT_GREEN),
        ("15s", "Person Click", "Graph modal\nshows network\nwith role-anchors",
         "User sees:\nWHICH POSITION\ndid what", ACCENT_PURPLE),
        ("30s", "Cross-Layer", "Another layer's\ncontradictions\nload in center",
         "Same person\nappears across\nmultiple layers", ACCENT_RED),
        ("5s", "Ask", "Natural language\nquery answered\nwith atom citations",
         "All layers for\nthis person\nin one response", ACCENT_YELLOW),
    ]

    for i, (time, title, content, insight, color) in enumerate(steps):
        x = Inches(0.4) + i * Inches(2.1)
        y = Inches(1.0)

        # Time badge
        badge = add_rect(slide, x + Inches(0.55), y, Inches(0.8), Inches(0.4), color)
        tf = badge.text_frame
        tf.paragraphs[0].text = time
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = TEXT_BRIGHT
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Step title
        add_text_box(slide, x, y + Inches(0.5), Inches(1.9), Inches(0.35),
                     title, 14, color, True, PP_ALIGN.CENTER)

        # Content card
        add_rect(slide, x, y + Inches(0.9), Inches(1.9), Inches(1.8), BG_CARD, color, Pt(1))
        add_multiline_text(slide, x + Inches(0.1), y + Inches(1.0), Inches(1.7), Inches(1.5),
                           [(line, 9, TEXT_WHITE) for line in content.split("\n")])

        # Insight card
        add_rect(slide, x, y + Inches(2.9), Inches(1.9), Inches(1.5), BG_PANEL)
        add_multiline_text(slide, x + Inches(0.1), y + Inches(3.0), Inches(1.7), Inches(1.3),
                           [(line, 9, TEXT_DIM) for line in insight.split("\n")])

        # Arrow
        if i < len(steps) - 1:
            add_text_box(slide, x + Inches(1.9), y + Inches(1.6), Inches(0.3), Inches(0.3),
                         ">", 18, TEXT_DIM, True, PP_ALIGN.CENTER)

    # Total time
    add_rect(slide, Inches(3.5), Inches(6.0), Inches(6.3), Inches(0.7), BG_CARD,
             ACCENT_GREEN, Pt(1))
    add_text_box(slide, Inches(3.8), Inches(6.1), Inches(5.7), Inches(0.5),
                 "Total: ~90 seconds to traverse 4 layers and grasp the proof structure",
                 14, ACCENT_GREEN, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════
# SLIDE 8 — Technical Architecture
# ════════════════════════════════════════
def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.45),
                 "Technical Architecture — 7 Layers", 24, ACCENT_BLUE, True)
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(12), Pt(1), ACCENT_BLUE)

    layers = [
        ("L1", "raw/", "Immutable Primary Sources", "13,495 records  |  PDFs  |  KakaoTalk  |  Scans", TEXT_DIM),
        ("L2", "wiki/", "Source of Truth", "Atoms + Hubs + Claims + Contradictions", ACCENT_GREEN),
        ("L3", "compiler/", "Deterministic Compiler", "atoms-to-cypher.py  |  build-record-index.py  |  rebuild-hubs.py", ACCENT_BLUE),
        ("L4", "Neo4j", "Materialized View (Cache)", "Graph = rebuilt from wiki anytime  |  Delete + recompile = no loss", ACCENT_YELLOW),
        ("L5", "search/", "Karpathy Retriever", "BM25 + LLM rerank + atom-citation  |  GraphRAG = auxiliary only", ACCENT_PURPLE),
        ("L6", "api/", "Single API", "/query  /atom/{id}  /neighborhood  /evidence  |  Pseudonym guard", LAYER_COLORS[5]),
        ("L7", "ui/", "Single-Screen UI", "Landing + 3-Panel Proof + Graph Modal  |  Angular + Cytoscape.js", ACCENT_BLUE),
    ]

    for i, (label, path, title, desc, color) in enumerate(layers):
        y = Inches(6.2) - i * Inches(0.78)
        w = Inches(11.5)

        # Layer bar
        add_rect(slide, Inches(0.9), y, w, Inches(0.65), BG_PANEL, color, Pt(1))

        # Label
        lbl = add_rect(slide, Inches(0.9), y, Inches(0.6), Inches(0.65), color)
        tf = lbl.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = TEXT_BRIGHT
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Path
        add_text_box(slide, Inches(1.7), y + Inches(0.05), Inches(1.2), Inches(0.3),
                     path, 11, color, True)
        # Title
        add_text_box(slide, Inches(3.0), y + Inches(0.05), Inches(3.0), Inches(0.3),
                     title, 12, TEXT_WHITE, True)
        # Description
        add_text_box(slide, Inches(1.7), y + Inches(0.33), Inches(9.5), Inches(0.25),
                     desc, 9, TEXT_DIM)

    # Arrow annotation
    add_text_box(slide, Inches(0.3), Inches(3.5), Inches(0.5), Inches(1.5),
                 "  ↑\n  |\n  |\ndata\nflow\n  |\n  ↓", 9, TEXT_DIM, False, PP_ALIGN.CENTER)

    # Invariant note
    add_text_box(slide, Inches(2.5), Inches(0.7), Inches(8.5), Inches(0.3),
                 "Invariant: raw/ → wiki/ → compiler → graph → API → UI  |  One-way data flow  |  Graph = cache  |  3x pseudonym guard",
                 10, ACCENT_GREEN, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════
# SLIDE 9 — Component Reuse from v1
# ════════════════════════════════════════
def slide_reuse(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.45),
                 "v1 Component Reuse Strategy", 24, ACCENT_BLUE, True)
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(12), Pt(1), ACCENT_BLUE)

    components = [
        ("graph.component.ts", "1,464", "REUSE", "Graph modal Cytoscape renderer", ACCENT_GREEN),
        ("evidence-browser.component.ts", "379", "REUSE", "Right panel evidence context", ACCENT_GREEN),
        ("pdf-viewer.component.ts", "—", "REUSE", "Scan preview in right panel", ACCENT_GREEN),
        ("search.component.ts", "536", "SIMPLIFY", "Reduce 5-mode → single input", ACCENT_YELLOW),
        ("command-center.component.ts", "692", "DECOMPOSE", "Extract useful parts only", ACCENT_YELLOW),
        ("layer-detail.component.ts", "—", "REWRITE", "Tower-style layer navigator", ACCENT_RED),
        ("dashboard.component.ts", "—", "REWRITE", "Proof declaration landing", ACCENT_RED),
        ("onboarding-tour.component.ts", "—", "REMOVE", "Landing is self-explanatory", TEXT_DIM),
    ]

    # Header
    add_rect(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.45), BG_CARD)
    headers = [("Component", 0.8, 3.5), ("LoC", 4.5, 0.8), ("Verdict", 5.5, 1.2),
               ("Role in v2", 6.9, 4.5)]
    for text, x, w in headers:
        add_text_box(slide, Inches(x), Inches(0.85), Inches(w), Inches(0.35),
                     text, 11, TEXT_DIM, True)

    for i, (name, loc, verdict, role, color) in enumerate(components):
        y = Inches(1.4) + i * Inches(0.6)

        if i % 2 == 0:
            add_rect(slide, Inches(0.6), y, Inches(12), Inches(0.5), BG_CARD)

        add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(3.5), Inches(0.35),
                     name, 11, TEXT_WHITE)
        add_text_box(slide, Inches(4.5), y + Inches(0.05), Inches(0.8), Inches(0.35),
                     loc, 11, TEXT_DIM, False, PP_ALIGN.CENTER)

        # Verdict badge
        badge = add_rect(slide, Inches(5.5), y + Inches(0.07), Inches(1.2), Inches(0.3), color)
        tf = badge.text_frame
        tf.paragraphs[0].text = verdict
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = TEXT_BRIGHT if color != TEXT_DIM else RGBColor(0x60, 0x60, 0x70)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_text_box(slide, Inches(6.9), y + Inches(0.05), Inches(5.5), Inches(0.35),
                     role, 11, TEXT_DIM)

    # Summary
    add_rect(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.8), BG_PANEL,
             ACCENT_BLUE, Pt(1))
    add_text_box(slide, Inches(0.9), Inches(6.4), Inches(11.4), Inches(0.6),
                 "3 REUSE  |  2 SIMPLIFY/DECOMPOSE  |  2 REWRITE  |  1 REMOVE\n"
                 "70% of v1's mature Cytoscape + Material code carries forward. Zero Cypher-form debt.",
                 12, TEXT_WHITE, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════
# MAIN
# ════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_title(prs)
    slide_paradigm(prs)
    slide_principles(prs)
    slide_landing(prs)
    slide_proof_screen(prs)
    slide_graph_modal(prs)
    slide_user_journey(prs)
    slide_architecture(prs)
    slide_reuse(prs)

    out_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "output")
    out_path = os.path.join(out_dir, "aurora-v2-ui-mockup-2026-04-11.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
