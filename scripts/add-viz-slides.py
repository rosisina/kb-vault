#!/usr/bin/env python3
"""
Add 5 visualization variant slides + comparison slide to the existing Aurora v2 UI mockup PPT.
Variants: Cytoscape.js (Aurora v1), Obsidian constellation, G6/AntV, 3D Force Graph, Sigma.js/Gephi
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os, math

# ── Shared palette (matches existing PPT) ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_PANEL = RGBColor(0x16, 0x21, 0x3E)
BG_CARD = RGBColor(0x1F, 0x2B, 0x4D)
BG_INPUT = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_BLUE = RGBColor(0x4E, 0x9A, 0xF5)
ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0xB0)
ACCENT_RED = RGBColor(0xE0, 0x6C, 0x75)
ACCENT_YELLOW = RGBColor(0xE5, 0xC0, 0x7B)
ACCENT_PURPLE = RGBColor(0xC6, 0x78, 0xDD)
ACCENT_ORANGE = RGBColor(0xFF, 0x92, 0x4C)
ACCENT_CYAN = RGBColor(0x56, 0xB6, 0xC2)
ACCENT_PINK = RGBColor(0xE0, 0x6C, 0xA5)
TEXT_WHITE = RGBColor(0xE0, 0xE0, 0xE0)
TEXT_DIM = RGBColor(0x8B, 0x8B, 0xA0)
TEXT_BRIGHT = RGBColor(0xFF, 0xFF, 0xFF)

# Aurora v1 layer colors (from graph.component.ts)
V1_LAYER_COLORS = {
    1: RGBColor(0xB7, 0x1C, 0x1C),  # Deep red
    2: RGBColor(0xE6, 0x51, 0x00),  # Deep orange
    3: RGBColor(0xF5, 0x7F, 0x17),  # Amber
    4: RGBColor(0x1B, 0x5E, 0x20),  # Deep green
    5: RGBColor(0x0D, 0x47, 0xA1),  # Deep blue
    6: RGBColor(0x4A, 0x14, 0x8C),  # Deep purple
    7: RGBColor(0x88, 0x0E, 0x4F),  # Deep pink
}

# Standard node data for all slides
NODES = [
    ("박성호", "원장-1", "Person", 4),
    ("이태호", "위원장-1", "Person", 4),
    ("이지영", "결재자-1", "Person", 3),
    ("이준호", "공모자-1", "Person", 4),
    ("장우진", "실무자-1", "Person", 2),
    ("KIATIS\n성능개선사업", "", "Event", 4),
    ("DIDC", "", "Org", 1),
    ("국전원", "", "Org", 3),
]


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.02
    return shape


def add_circle(slide, cx, cy, r, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_diamond(slide, cx, cy, size, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, cx - size//2, cy - size//2, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Pretendard"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, default_size=11,
                       default_color=TEXT_WHITE, font_name="Pretendard"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(2)
    return txBox


def add_line_connector(slide, x1, y1, x2, y2, color, width=Pt(1)):
    """Add a line between two points."""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def slide_section_header(prs):
    """Section divider slide for visualization variants."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_rect(slide, Inches(2), Inches(2.8), Inches(9.333), Pt(2), ACCENT_BLUE)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11.333), Inches(1),
                 "Relationship Network Visualization", 36, ACCENT_BLUE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11.333), Inches(0.5),
                 "5 Visualization Approaches for Layer 4 Evidence Network", 18, TEXT_DIM, False, PP_ALIGN.CENTER)

    add_multiline_text(slide, Inches(2), Inches(3.3), Inches(9.333), Inches(3.5), [
        ("Variant A    Cytoscape.js — Aurora v1 Style (현재 구축 완료)", 14, ACCENT_GREEN, True),
        ("Variant B    Obsidian Constellation — Knowledge Graph Style", 14, ACCENT_PURPLE, True),
        ("Variant C    G6/AntV — Enterprise Dark Theme", 14, ACCENT_ORANGE, True),
        ("Variant D    3D Force Graph — Three.js Immersive", 14, ACCENT_CYAN, True),
        ("Variant E    Sigma.js/Gephi — Large-Scale Network Science", 14, ACCENT_YELLOW, True),
        ("", 10, TEXT_DIM),
        ("Each variant shows the same Layer 4 evidence network", 12, TEXT_DIM),
        ("with identical data: 5 persons, 1 event, 2 organizations", 12, TEXT_DIM),
    ])


# ════════════════════════════════════════
# VARIANT A — Cytoscape.js (Aurora v1 actual style)
# ════════════════════════════════════════
def slide_variant_a_cytoscape(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0xFF, 0xFF, 0xFF))  # v1 uses light background

    add_text_box(slide, Inches(0.4), Inches(0.15), Inches(9), Inches(0.45),
                 "Variant A — Cytoscape.js (Aurora v1 Actual Style)", 24, RGBColor(0x33, 0x33, 0x33), True)
    add_text_box(slide, Inches(0.4), Inches(0.55), Inches(9), Inches(0.3),
                 "graph.component.ts  |  1,464 LoC  |  COSE force-directed  |  Canvas 2D", 11,
                 RGBColor(0x99, 0x99, 0x99))

    # Tag
    tag = add_rect(slide, Inches(10.5), Inches(0.2), Inches(2.5), Inches(0.4), ACCENT_GREEN)
    tf = tag.text_frame
    tf.paragraphs[0].text = "ALREADY BUILT"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = TEXT_BRIGHT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Graph area — light background (v1 style)
    graph_bg = RGBColor(0xFA, 0xFA, 0xFA)
    add_rect(slide, Inches(0.4), Inches(0.9), Inches(9.0), Inches(5.8), graph_bg,
             RGBColor(0xE0, 0xE0, 0xE0), Pt(1))

    # v1 node positions (COSE layout simulation)
    person_nodes = [
        (Inches(2.5), Inches(2.0), "박성호", "(원장-1)", V1_LAYER_COLORS[4], RGBColor(0xC6, 0x28, 0x28), Inches(0.45)),
        (Inches(6.0), Inches(1.5), "이태호", "(위원장-1)", V1_LAYER_COLORS[4], RGBColor(0xC6, 0x28, 0x28), Inches(0.40)),
        (Inches(1.5), Inches(4.2), "이지영", "(결재자-1)", V1_LAYER_COLORS[3], RGBColor(0xC6, 0x28, 0x28), Inches(0.42)),
        (Inches(5.5), Inches(4.5), "이준호", "(공모자-1)", V1_LAYER_COLORS[4], RGBColor(0xC6, 0x28, 0x28), Inches(0.38)),
        (Inches(3.5), Inches(5.5), "장우진", "(실무자-1)", V1_LAYER_COLORS[2], RGBColor(0x15, 0x65, 0xC0), Inches(0.35)),
    ]

    for x, y, name, role, border_color, fill_color, radius in person_nodes:
        # Circle node with layer-colored border (v1 style)
        node = add_circle(slide, x, y, radius, fill_color, border_color, Pt(3))
        add_text_box(slide, x - Inches(0.5), y + radius + Pt(2), Inches(1.0), Inches(0.25),
                     name, 8, RGBColor(0x33, 0x33, 0x33), False, PP_ALIGN.CENTER)

    # Event node — round-diamond (v1 style for CriminalAct)
    ev_x, ev_y = Inches(4.0), Inches(3.0)
    ev_shape = add_diamond(slide, ev_x, ev_y, Inches(0.7), V1_LAYER_COLORS[4], RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide, ev_x - Inches(0.7), ev_y + Inches(0.45), Inches(1.4), Inches(0.4),
                 "KIATIS\n성능개선사업", 7, RGBColor(0x33, 0x33, 0x33), False, PP_ALIGN.CENTER)

    # Org nodes — rectangles (v1 style)
    org_nodes = [
        (Inches(7.5), Inches(3.0), "DIDC", V1_LAYER_COLORS[1]),
        (Inches(7.5), Inches(5.0), "국전원", V1_LAYER_COLORS[3]),
    ]
    for x, y, name, color in org_nodes:
        org = add_rect(slide, x - Inches(0.45), y - Inches(0.25), Inches(0.9), Inches(0.5),
                       RGBColor(0x37, 0x47, 0x4F), RGBColor(0xB0, 0xBE, 0xC5), Pt(3))
        tf = org.text_frame
        tf.paragraphs[0].text = name
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = TEXT_BRIGHT
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Edge labels
    edge_labels = [
        (Inches(3.0), Inches(2.3), "ORDERED", RGBColor(0xC6, 0x28, 0x28), 7),
        (Inches(5.0), Inches(2.0), "시험평가", RGBColor(0xE6, 0x51, 0x00), 7),
        (Inches(1.8), Inches(3.2), "EXECUTED", RGBColor(0xE6, 0x51, 0x00), 7),
        (Inches(3.5), Inches(4.8), "VICTIM", RGBColor(0x15, 0x65, 0xC0), 7),
        (Inches(4.5), Inches(4.0), "공모", RGBColor(0xC6, 0x28, 0x28), 7),
        (Inches(6.0), Inches(3.0), "AFFILIATED", RGBColor(0x90, 0xA4, 0xAE), 7),
    ]
    for x, y, label, color, size in edge_labels:
        add_text_box(slide, x, y, Inches(1.0), Inches(0.2), f"--{label}-->", size, color, False, PP_ALIGN.CENTER)

    # Right panel — info panel (v1 style)
    add_rect(slide, Inches(9.6), Inches(0.9), Inches(3.4), Inches(5.8),
             RGBColor(0xF5, 0xF5, 0xF5), RGBColor(0xE0, 0xE0, 0xE0), Pt(1))

    add_multiline_text(slide, Inches(9.8), Inches(1.0), Inches(3.0), Inches(5.5), [
        ("Node Information", 14, RGBColor(0x33, 0x33, 0x33), True),
        ("", 6, TEXT_DIM),
        ("Type: Person (Perpetrator)", 10, RGBColor(0x66, 0x66, 0x66)),
        ("Layer: 4", 10, RGBColor(0x66, 0x66, 0x66)),
        ("Name: 박성호 (원장-1)", 10, RGBColor(0x33, 0x33, 0x33), True),
        ("", 6, TEXT_DIM),
        ("Connected Evidence:", 10, RGBColor(0x33, 0x33, 0x33), True),
        ("  Record No. 10,347", 9, RGBColor(0x0D, 0x47, 0xA1)),
        ("  Record No. 10,352", 9, RGBColor(0x0D, 0x47, 0xA1)),
        ("  Record No. 10,389", 9, RGBColor(0x0D, 0x47, 0xA1)),
        ("", 6, TEXT_DIM),
        ("Criminal Acts:", 10, RGBColor(0x33, 0x33, 0x33), True),
        ("  시험평가 조작 (L4)", 9, V1_LAYER_COLORS[4]),
        ("  사업관리 은폐 (L3)", 9, V1_LAYER_COLORS[3]),
        ("", 6, TEXT_DIM),
        ("Characteristics:", 10, RGBColor(0x33, 0x33, 0x33), True),
        ("  Light background", 9, RGBColor(0x99, 0x99, 0x99)),
        ("  Circle = Person", 9, RGBColor(0x99, 0x99, 0x99)),
        ("  Rectangle = Organization", 9, RGBColor(0x99, 0x99, 0x99)),
        ("  Diamond = Criminal Act", 9, RGBColor(0x99, 0x99, 0x99)),
        ("  Border color = Layer", 9, RGBColor(0x99, 0x99, 0x99)),
        ("  Fill color = Side", 9, RGBColor(0x99, 0x99, 0x99)),
        ("    (Red=perp, Blue=victim)", 9, RGBColor(0x99, 0x99, 0x99)),
    ])


# ════════════════════════════════════════
# VARIANT B — Obsidian Constellation Style
# ════════════════════════════════════════
def slide_variant_b_obsidian(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x0A, 0x0A, 0x14))  # Very dark (Obsidian-style)

    add_text_box(slide, Inches(0.4), Inches(0.15), Inches(9), Inches(0.45),
                 "Variant B — Obsidian Constellation Style", 24, ACCENT_PURPLE, True)
    add_text_box(slide, Inches(0.4), Inches(0.55), Inches(9), Inches(0.3),
                 "PIXI.js (WebGL)  |  Minimal nodes  |  Labels on hover  |  Galaxy aesthetic", 11, TEXT_DIM)

    # Tag
    tag = add_rect(slide, Inches(10.5), Inches(0.2), Inches(2.5), Inches(0.4), ACCENT_PURPLE)
    tf = tag.text_frame
    tf.paragraphs[0].text = "INSPIRATION"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = TEXT_BRIGHT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Graph area — very dark
    graph_bg = RGBColor(0x08, 0x08, 0x12)
    add_rect(slide, Inches(0.4), Inches(0.9), Inches(9.0), Inches(5.8), graph_bg)

    # Constellation dots — small glowing nodes
    # Obsidian style: tiny bright dots, thin translucent lines, labels only on hover
    import random
    random.seed(42)

    # Background constellation dots (ambient nodes — other vault notes)
    ambient_colors = [
        RGBColor(0x40, 0x50, 0x70),
        RGBColor(0x35, 0x45, 0x65),
        RGBColor(0x30, 0x40, 0x55),
    ]
    for _ in range(45):
        ax = Inches(0.6 + random.random() * 8.6)
        ay = Inches(1.1 + random.random() * 5.4)
        ar = Inches(0.03 + random.random() * 0.04)
        ac = random.choice(ambient_colors)
        add_circle(slide, ax, ay, ar, ac)

    # Thin connection lines (ambient — simulated as very small text)
    for _ in range(20):
        lx = Inches(0.8 + random.random() * 8.2)
        ly = Inches(1.3 + random.random() * 5.0)
        add_text_box(slide, lx, ly, Inches(0.8), Inches(0.02),
                     "—" * 8, 3, RGBColor(0x20, 0x28, 0x40))

    # Primary constellation nodes (Layer 4 evidence network)
    primary_nodes = [
        (Inches(3.0), Inches(2.5), "박성호", Inches(0.12), RGBColor(0xE0, 0x6C, 0x75)),
        (Inches(6.5), Inches(2.0), "이태호", Inches(0.10), RGBColor(0xE0, 0x6C, 0x75)),
        (Inches(2.0), Inches(4.5), "이지영", Inches(0.11), RGBColor(0xE5, 0xC0, 0x7B)),
        (Inches(6.0), Inches(4.8), "이준호", Inches(0.10), RGBColor(0xE0, 0x6C, 0x75)),
        (Inches(4.0), Inches(5.5), "장우진", Inches(0.08), ACCENT_BLUE),
        (Inches(4.5), Inches(3.3), "KIATIS", Inches(0.14), ACCENT_PURPLE),
        (Inches(7.8), Inches(3.0), "DIDC", Inches(0.13), ACCENT_CYAN),
        (Inches(7.5), Inches(5.0), "국전원", Inches(0.11), ACCENT_GREEN),
    ]

    # Glow halos (outer ring — simulates Obsidian glow)
    for x, y, name, r, color in primary_nodes:
        # Outer glow (larger, transparent)
        glow_r = r + Inches(0.08)
        glow_color = RGBColor(
            min(255, color.red + 30) if hasattr(color, 'red') else 0x80,
            min(255, color.green + 30) if hasattr(color, 'green') else 0x80,
            min(255, color.blue + 30) if hasattr(color, 'blue') else 0x80,
        )
        # Simulating glow with a slightly larger, dimmer circle
        dim_color = RGBColor(color[0] // 3, color[1] // 3, color[2] // 3)
        add_circle(slide, x, y, glow_r, dim_color)
        # Core dot
        add_circle(slide, x, y, r, color)

    # Labels — only shown for hovered node (simulate one hovered)
    # Show label for "박성호" as if hovered
    add_text_box(slide, Inches(2.3), Inches(2.15), Inches(1.4), Inches(0.25),
                 "박성호 (원장-1)", 10, TEXT_BRIGHT, False, PP_ALIGN.CENTER)
    # Other labels very dim (as if not hovered)
    for x, y, name, r, color in primary_nodes[1:]:
        add_text_box(slide, x - Inches(0.5), y + r + Pt(3), Inches(1.0), Inches(0.2),
                     name, 7, RGBColor(0x40, 0x50, 0x65), False, PP_ALIGN.CENTER)

    # Connection lines between primary nodes (thin, translucent)
    connections = [
        (Inches(3.0), Inches(2.5), Inches(4.5), Inches(3.3), RGBColor(0x50, 0x30, 0x50)),
        (Inches(6.5), Inches(2.0), Inches(4.5), Inches(3.3), RGBColor(0x50, 0x30, 0x50)),
        (Inches(2.0), Inches(4.5), Inches(4.5), Inches(3.3), RGBColor(0x40, 0x40, 0x30)),
        (Inches(6.0), Inches(4.8), Inches(4.5), Inches(3.3), RGBColor(0x50, 0x30, 0x50)),
        (Inches(3.0), Inches(2.5), Inches(2.0), Inches(4.5), RGBColor(0x30, 0x30, 0x50)),
        (Inches(6.0), Inches(4.8), Inches(2.0), Inches(4.5), RGBColor(0x50, 0x30, 0x30)),
        (Inches(4.5), Inches(3.3), Inches(7.8), Inches(3.0), RGBColor(0x30, 0x40, 0x50)),
        (Inches(4.5), Inches(3.3), Inches(7.5), Inches(5.0), RGBColor(0x30, 0x50, 0x40)),
    ]
    for x1, y1, x2, y2, color in connections:
        try:
            add_line_connector(slide, x1, y1, x2, y2, color, Pt(0.75))
        except:
            pass  # Connector API may not support all positions

    # Right panel — Obsidian-style local graph info
    add_rect(slide, Inches(9.6), Inches(0.9), Inches(3.4), Inches(5.8),
             RGBColor(0x10, 0x14, 0x22), RGBColor(0x25, 0x30, 0x45), Pt(1))

    add_multiline_text(slide, Inches(9.8), Inches(1.0), Inches(3.0), Inches(5.5), [
        ("Local Graph View", 14, ACCENT_PURPLE, True),
        ("Focus: 박성호 (원장-1)", 11, TEXT_WHITE, True),
        ("", 6, TEXT_DIM),
        ("Connections: 4", 10, TEXT_DIM),
        ("  KIATIS 성능개선사업", 9, ACCENT_PURPLE),
        ("  이지영 (지시)", 9, ACCENT_YELLOW),
        ("  이태호 (시험평가)", 9, ACCENT_RED),
        ("  DIDC (소속)", 9, ACCENT_CYAN),
        ("", 6, TEXT_DIM),
        ("Backlinks: 6 atoms", 10, TEXT_DIM),
        ("  prosecution-misapplies...", 8, ACCENT_BLUE),
        ("  layer4-evaluation-...", 8, ACCENT_BLUE),
        ("  2436ho-didc-naming...", 8, ACCENT_BLUE),
        ("", 10, TEXT_DIM),
        ("─────────────────────", 8, RGBColor(0x25, 0x30, 0x45)),
        ("", 6, TEXT_DIM),
        ("Characteristics:", 10, TEXT_WHITE, True),
        ("  Very dark background", 9, TEXT_DIM),
        ("  Tiny bright dots = notes", 9, TEXT_DIM),
        ("  Label on hover only", 9, TEXT_DIM),
        ("  Glow halo = connectivity", 9, TEXT_DIM),
        ("  Ambient dots = full vault", 9, TEXT_DIM),
        ("  Galaxy / constellation feel", 9, TEXT_DIM),
    ])


# ════════════════════════════════════════
# VARIANT C — G6/AntV Enterprise Dark Theme
# ════════════════════════════════════════
def slide_variant_c_g6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x14, 0x17, 0x1F))  # G6 dark theme bg

    add_text_box(slide, Inches(0.4), Inches(0.15), Inches(9), Inches(0.45),
                 "Variant C — G6/AntV Enterprise Dark Theme", 24, ACCENT_ORANGE, True)
    add_text_box(slide, Inches(0.4), Inches(0.55), Inches(9), Inches(0.3),
                 "Canvas/WebGL  |  Ant Design system  |  Built-in dark theme  |  Rich node content", 11, TEXT_DIM)

    # Tag
    tag = add_rect(slide, Inches(10.5), Inches(0.2), Inches(2.5), Inches(0.4), ACCENT_ORANGE)
    tf = tag.text_frame
    tf.paragraphs[0].text = "RECOMMENDED"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = TEXT_BRIGHT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Graph area
    graph_bg = RGBColor(0x11, 0x14, 0x1C)
    add_rect(slide, Inches(0.4), Inches(0.9), Inches(9.0), Inches(5.8), graph_bg)

    # G6 style: rich node cards with multiple info lines
    # Person nodes as rich cards
    g6_nodes = [
        (Inches(1.8), Inches(1.8), "박성호", "원장-1", "L4 Perpetrator", ACCENT_RED, RGBColor(0x2A, 0x15, 0x15)),
        (Inches(5.8), Inches(1.3), "이태호", "위원장-1", "L4 Perpetrator", ACCENT_RED, RGBColor(0x2A, 0x15, 0x15)),
        (Inches(1.2), Inches(4.5), "이지영", "결재자-1", "L3 Perpetrator", ACCENT_YELLOW, RGBColor(0x2A, 0x25, 0x12)),
        (Inches(6.2), Inches(4.8), "이준호", "공모자-1", "L4 Perpetrator", ACCENT_RED, RGBColor(0x2A, 0x15, 0x15)),
        (Inches(3.5), Inches(5.8), "장우진", "실무자-1", "L2 Victim", ACCENT_BLUE, RGBColor(0x12, 0x18, 0x2A)),
    ]

    for x, y, name, role, layer_tag, accent, card_bg in g6_nodes:
        card_w, card_h = Inches(1.5), Inches(1.0)
        # Card background
        card = add_rect(slide, x, y, card_w, card_h, card_bg, accent, Pt(1.5))
        # Status dot
        dot = add_circle(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.06), accent)
        # Name
        add_text_box(slide, x + Inches(0.25), y + Inches(0.05), Inches(1.15), Inches(0.25),
                     name, 11, TEXT_BRIGHT, True)
        # Role
        add_text_box(slide, x + Inches(0.1), y + Inches(0.35), Inches(1.3), Inches(0.2),
                     f"({role})", 9, TEXT_DIM)
        # Layer tag
        tag_shape = add_rect(slide, x + Inches(0.1), y + Inches(0.6), Inches(1.1), Inches(0.25), accent)
        tf2 = tag_shape.text_frame
        tf2.paragraphs[0].text = layer_tag
        tf2.paragraphs[0].font.size = Pt(7)
        tf2.paragraphs[0].font.color.rgb = TEXT_BRIGHT
        tf2.paragraphs[0].font.bold = True
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Event node — special card
    ev_x, ev_y = Inches(3.8), Inches(3.0)
    ev_card = add_rect(slide, ev_x, ev_y, Inches(1.8), Inches(1.2),
                       RGBColor(0x1A, 0x12, 0x2A), ACCENT_PURPLE, Pt(2))
    add_text_box(slide, ev_x + Inches(0.1), ev_y + Inches(0.1), Inches(1.6), Inches(0.3),
                 "KIATIS 성능개선사업", 10, ACCENT_PURPLE, True)
    add_text_box(slide, ev_x + Inches(0.1), ev_y + Inches(0.4), Inches(1.6), Inches(0.2),
                 "2018-2019  |  6.25억", 8, TEXT_DIM)
    add_text_box(slide, ev_x + Inches(0.1), ev_y + Inches(0.65), Inches(1.6), Inches(0.2),
                 "Layer 4 Core Event", 8, ACCENT_PURPLE)
    # Animated flow indicator (dots)
    add_text_box(slide, ev_x + Inches(0.1), ev_y + Inches(0.9), Inches(1.6), Inches(0.2),
                 ">>> flow >>>", 7, RGBColor(0x40, 0x30, 0x55))

    # Org nodes — combo group style
    org_nodes = [
        (Inches(7.3), Inches(2.5), "DIDC", "국방통합데이터센터", ACCENT_CYAN),
        (Inches(7.5), Inches(4.5), "국전원", "국방전산원", ACCENT_GREEN),
    ]
    for x, y, name, full, color in org_nodes:
        add_rect(slide, x, y, Inches(1.5), Inches(0.8),
                 RGBColor(0x12, 0x1C, 0x20), color, Pt(1.5))
        add_text_box(slide, x + Inches(0.1), y + Inches(0.05), Inches(1.3), Inches(0.25),
                     name, 12, color, True)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.35), Inches(1.3), Inches(0.2),
                     full, 7, TEXT_DIM)
        # Member count badge
        add_rect(slide, x + Inches(1.1), y + Inches(0.5), Inches(0.3), Inches(0.2), color)
        add_text_box(slide, x + Inches(1.1), y + Inches(0.5), Inches(0.3), Inches(0.2),
                     "3", 8, TEXT_BRIGHT, True, PP_ALIGN.CENTER)

    # Edge labels with animated flow dots
    edges = [
        (Inches(3.0), Inches(2.3), "ORDERED_BY", ACCENT_RED),
        (Inches(5.2), Inches(1.8), "시험평가지시", ACCENT_ORANGE),
        (Inches(2.0), Inches(3.5), "EXECUTED", ACCENT_ORANGE),
        (Inches(4.5), Inches(5.0), "공모", ACCENT_RED),
        (Inches(3.0), Inches(5.5), "VICTIMIZED", ACCENT_BLUE),
        (Inches(6.5), Inches(3.0), "소속", ACCENT_CYAN),
    ]
    for x, y, label, color in edges:
        add_text_box(slide, x, y, Inches(1.2), Inches(0.2),
                     f"-- {label} -->", 7, color, False, PP_ALIGN.CENTER)

    # Right panel — G6 style info
    add_rect(slide, Inches(9.6), Inches(0.9), Inches(3.4), Inches(5.8),
             RGBColor(0x14, 0x17, 0x1F), RGBColor(0x30, 0x35, 0x45), Pt(1))

    add_multiline_text(slide, Inches(9.8), Inches(1.0), Inches(3.0), Inches(5.5), [
        ("G6 v5 — AntV", 14, ACCENT_ORANGE, True),
        ("Ant Design Dark Theme", 11, TEXT_DIM),
        ("", 6, TEXT_DIM),
        ("Why Recommended:", 11, TEXT_WHITE, True),
        ("", 4, TEXT_DIM),
        ("  Built-in dark theme", 9, ACCENT_GREEN),
        ("  (no custom CSS needed)", 8, TEXT_DIM),
        ("", 4, TEXT_DIM),
        ("  Rich node cards", 9, ACCENT_GREEN),
        ("  (role, layer, status inline)", 8, TEXT_DIM),
        ("", 4, TEXT_DIM),
        ("  Combo grouping", 9, ACCENT_GREEN),
        ("  (group by Layer 1-7)", 8, TEXT_DIM),
        ("", 4, TEXT_DIM),
        ("  Animated edge flow", 9, ACCENT_GREEN),
        ("  (shows causality direction)", 8, TEXT_DIM),
        ("", 4, TEXT_DIM),
        ("  WASM-powered layouts", 9, ACCENT_GREEN),
        ("  (Rust compiled, fast)", 8, TEXT_DIM),
        ("", 6, TEXT_DIM),
        ("─────────────────────", 8, RGBColor(0x30, 0x35, 0x45)),
        ("Angular: vanilla JS/TS", 9, TEXT_DIM),
        ("  integration via ElementRef", 8, TEXT_DIM),
        ("Performance: ~10K nodes", 9, TEXT_DIM),
    ])


# ════════════════════════════════════════
# VARIANT D — 3D Force Graph (Three.js)
# ════════════════════════════════════════
def slide_variant_d_3d(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x02, 0x02, 0x08))  # Near-black (space)

    add_text_box(slide, Inches(0.4), Inches(0.15), Inches(9), Inches(0.45),
                 "Variant D — 3D Force Graph (Three.js)", 24, ACCENT_CYAN, True)
    add_text_box(slide, Inches(0.4), Inches(0.55), Inches(9), Inches(0.3),
                 "WebGL  |  Three.js  |  3d-force-graph  |  Orbit camera  |  Bloom post-processing", 11, TEXT_DIM)

    # Tag
    tag = add_rect(slide, Inches(10.5), Inches(0.2), Inches(2.5), Inches(0.4), ACCENT_CYAN)
    tf = tag.text_frame
    tf.paragraphs[0].text = "WOW FACTOR"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = TEXT_BRIGHT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Graph area — deep space black
    add_rect(slide, Inches(0.4), Inches(0.9), Inches(9.0), Inches(5.8), RGBColor(0x01, 0x01, 0x05))

    # 3D perspective simulation — nodes at various "depths"
    # Foreground (large, bright)
    fg_nodes = [
        (Inches(3.5), Inches(2.8), "박성호", "(원장-1)", Inches(0.25), ACCENT_RED, True),
        (Inches(5.0), Inches(3.5), "KIATIS", "", Inches(0.30), ACCENT_PURPLE, True),
    ]
    # Midground (medium)
    mg_nodes = [
        (Inches(6.8), Inches(2.0), "이태호", "(위원장-1)", Inches(0.17), ACCENT_RED, False),
        (Inches(2.0), Inches(4.5), "이지영", "(결재자-1)", Inches(0.18), ACCENT_YELLOW, False),
        (Inches(6.5), Inches(5.0), "이준호", "(공모자-1)", Inches(0.16), ACCENT_RED, False),
    ]
    # Background (small, dim)
    bg_nodes = [
        (Inches(4.0), Inches(5.8), "장우진", "", Inches(0.10), ACCENT_BLUE, False),
        (Inches(8.0), Inches(3.5), "DIDC", "", Inches(0.12), ACCENT_CYAN, False),
        (Inches(7.5), Inches(5.5), "국전원", "", Inches(0.11), ACCENT_GREEN, False),
    ]

    # Draw depth-sorted (background first)
    for x, y, name, role, r, color, is_focus in bg_nodes:
        dim = RGBColor(color[0] // 3, color[1] // 3, color[2] // 3)
        add_circle(slide, x, y, r + Inches(0.04), dim)  # Glow
        add_circle(slide, x, y, r, RGBColor(color[0] // 2, color[1] // 2, color[2] // 2))
        add_text_box(slide, x - Inches(0.4), y + r + Pt(2), Inches(0.8), Inches(0.2),
                     name, 7, RGBColor(0x50, 0x50, 0x60), False, PP_ALIGN.CENTER)

    for x, y, name, role, r, color, is_focus in mg_nodes:
        dim = RGBColor(color[0] // 4, color[1] // 4, color[2] // 4)
        add_circle(slide, x, y, r + Inches(0.06), dim)
        add_circle(slide, x, y, r, color)
        label = f"{name} {role}" if role else name
        add_text_box(slide, x - Inches(0.6), y + r + Pt(3), Inches(1.2), Inches(0.2),
                     label, 8, TEXT_DIM, False, PP_ALIGN.CENTER)

    for x, y, name, role, r, color, is_focus in fg_nodes:
        # Bloom glow effect (multiple concentric circles)
        for g in range(3, 0, -1):
            glow_r = r + Inches(0.05 * g)
            glow_c = RGBColor(color[0] // (g + 1), color[1] // (g + 1), color[2] // (g + 1))
            add_circle(slide, x, y, glow_r, glow_c)
        add_circle(slide, x, y, r, color)
        label = f"{name} {role}" if role else name
        add_text_box(slide, x - Inches(0.6), y + r + Pt(4), Inches(1.2), Inches(0.3),
                     label, 10, TEXT_BRIGHT, True, PP_ALIGN.CENTER)

    # 3D perspective lines (edges receding into "space")
    edge_hints = [
        (Inches(3.7), Inches(3.0), Inches(5.0), Inches(3.5), RGBColor(0x40, 0x20, 0x40)),
        (Inches(5.2), Inches(3.5), Inches(6.8), Inches(2.0), RGBColor(0x30, 0x20, 0x30)),
        (Inches(3.3), Inches(2.9), Inches(2.0), Inches(4.5), RGBColor(0x30, 0x30, 0x40)),
        (Inches(5.2), Inches(3.6), Inches(6.5), Inches(5.0), RGBColor(0x30, 0x20, 0x30)),
        (Inches(5.1), Inches(3.6), Inches(8.0), Inches(3.5), RGBColor(0x20, 0x30, 0x35)),
    ]
    for x1, y1, x2, y2, color in edge_hints:
        try:
            add_line_connector(slide, x1, y1, x2, y2, color, Pt(0.5))
        except:
            pass

    # Camera orbit hint
    add_text_box(slide, Inches(1.0), Inches(6.2), Inches(3.0), Inches(0.3),
                 "  Drag to orbit  |  Scroll to zoom  |  Click to fly-to", 9, TEXT_DIM)

    # Right panel
    add_rect(slide, Inches(9.6), Inches(0.9), Inches(3.4), Inches(5.8),
             RGBColor(0x05, 0x05, 0x10), RGBColor(0x20, 0x30, 0x45), Pt(1))

    add_multiline_text(slide, Inches(9.8), Inches(1.0), Inches(3.0), Inches(5.5), [
        ("3D Force Graph", 14, ACCENT_CYAN, True),
        ("Three.js (WebGL)", 11, TEXT_DIM),
        ("", 6, TEXT_DIM),
        ("Visual Impact:", 11, TEXT_WHITE, True),
        ("  Highest of all variants", 9, ACCENT_GREEN),
        ("  Nodes float in 3D space", 9, TEXT_DIM),
        ("  Bloom post-processing", 9, TEXT_DIM),
        ("  Depth = perceived importance", 9, TEXT_DIM),
        ("", 6, TEXT_DIM),
        ("Interaction:", 11, TEXT_WHITE, True),
        ("  Orbit camera (rotate)", 9, TEXT_DIM),
        ("  Click = fly-to-node", 9, TEXT_DIM),
        ("  Hover = show connections", 9, TEXT_DIM),
        ("  Continuous physics", 9, TEXT_DIM),
        ("", 6, TEXT_DIM),
        ("Trade-offs:", 11, TEXT_WHITE, True),
        ("  3D adds cognitive load", 9, ACCENT_YELLOW),
        ("  Not ideal for daily analysis", 9, ACCENT_YELLOW),
        ("  Best for presentations", 9, ACCENT_GREEN),
        ("  Angular: use vanilla lib", 9, TEXT_DIM),
        ("  (3d-force-graph, not React)", 9, TEXT_DIM),
        ("", 6, TEXT_DIM),
        ("Use case:", 11, TEXT_WHITE, True),
        ("  Hero visualization", 9, ACCENT_CYAN),
        ("  Court presentations", 9, ACCENT_CYAN),
        ("  Stakeholder demos", 9, ACCENT_CYAN),
    ])


# ════════════════════════════════════════
# VARIANT E — Sigma.js / Gephi (Network Science)
# ════════════════════════════════════════
def slide_variant_e_sigma(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x12, 0x12, 0x18))

    add_text_box(slide, Inches(0.4), Inches(0.15), Inches(9), Inches(0.45),
                 "Variant E — Sigma.js / Gephi Pipeline", 24, ACCENT_YELLOW, True)
    add_text_box(slide, Inches(0.4), Inches(0.55), Inches(9), Inches(0.3),
                 "WebGL  |  ForceAtlas2 layout  |  100K+ nodes  |  Gephi → Sigma pipeline", 11, TEXT_DIM)

    # Tag
    tag = add_rect(slide, Inches(10.5), Inches(0.2), Inches(2.5), Inches(0.4), ACCENT_YELLOW)
    tf = tag.text_frame
    tf.paragraphs[0].text = "SCALE CHAMPION"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = RGBColor(0x10, 0x10, 0x10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Graph area
    add_rect(slide, Inches(0.4), Inches(0.9), Inches(9.0), Inches(5.8), RGBColor(0x0D, 0x0D, 0x12))

    # Sigma/Gephi style: ForceAtlas2 produces organic clustering
    # Simulate community detection with dense clusters
    import random
    random.seed(99)

    # Community clusters (7 layers as communities)
    communities = [
        (Inches(2.5), Inches(2.5), ACCENT_RED, 12, "L4"),       # Layer 4 cluster (center)
        (Inches(5.5), Inches(2.0), ACCENT_YELLOW, 8, "L3"),     # Layer 3
        (Inches(7.0), Inches(3.5), ACCENT_BLUE, 10, "L1"),      # Layer 1
        (Inches(2.0), Inches(5.0), ACCENT_PURPLE, 7, "L5"),     # Layer 5
        (Inches(5.0), Inches(5.5), ACCENT_ORANGE, 9, "L6"),     # Layer 6
        (Inches(7.5), Inches(5.5), ACCENT_CYAN, 6, "L7"),       # Layer 7
        (Inches(4.0), Inches(1.5), ACCENT_GREEN, 5, "L2"),      # Layer 2
    ]

    # Draw community dots (many small dots in clusters)
    for cx, cy, color, count, label in communities:
        dim_color = RGBColor(color[0] // 3, color[1] // 3, color[2] // 3)
        for _ in range(count):
            dx = Inches((random.random() - 0.5) * 1.2)
            dy = Inches((random.random() - 0.5) * 1.2)
            r = Inches(0.03 + random.random() * 0.05)
            add_circle(slide, cx + dx, cy + dy, r, color)
        # Inter-cluster edges (thin lines)
        for _ in range(count // 2):
            dx = Inches((random.random() - 0.5) * 1.0)
            dy = Inches((random.random() - 0.5) * 1.0)
            add_text_box(slide, cx + dx - Inches(0.3), cy + dy, Inches(0.6), Inches(0.01),
                         "—" * 5, 2, dim_color)
        # Community label
        add_text_box(slide, cx - Inches(0.3), cy - Inches(0.8), Inches(0.6), Inches(0.3),
                     label, 14, color, True, PP_ALIGN.CENTER)

    # Named nodes (larger, in their community positions)
    named = [
        (Inches(2.3), Inches(2.3), "박성호", ACCENT_RED, Inches(0.1)),
        (Inches(2.8), Inches(2.7), "이준호", ACCENT_RED, Inches(0.08)),
        (Inches(5.3), Inches(1.8), "이지영", ACCENT_YELLOW, Inches(0.09)),
        (Inches(7.2), Inches(3.3), "DIDC", ACCENT_BLUE, Inches(0.12)),
        (Inches(5.5), Inches(2.3), "국전원", ACCENT_YELLOW, Inches(0.1)),
        (Inches(2.5), Inches(3.0), "KIATIS", ACCENT_RED, Inches(0.13)),
    ]
    for x, y, name, color, r in named:
        add_circle(slide, x, y, r + Inches(0.03), RGBColor(color[0] // 4, color[1] // 4, color[2] // 4))
        add_circle(slide, x, y, r, color)
        add_text_box(slide, x - Inches(0.4), y + r + Pt(2), Inches(0.8), Inches(0.2),
                     name, 7, TEXT_WHITE, False, PP_ALIGN.CENTER)

    # Cross-community edges
    cross_edges = [
        (Inches(2.5), Inches(2.5), Inches(5.5), Inches(2.0), RGBColor(0x50, 0x40, 0x20)),
        (Inches(2.5), Inches(2.5), Inches(7.0), Inches(3.5), RGBColor(0x40, 0x20, 0x20)),
        (Inches(5.5), Inches(2.0), Inches(7.0), Inches(3.5), RGBColor(0x30, 0x40, 0x30)),
        (Inches(2.5), Inches(2.5), Inches(2.0), Inches(5.0), RGBColor(0x40, 0x20, 0x40)),
        (Inches(2.5), Inches(2.5), Inches(5.0), Inches(5.5), RGBColor(0x50, 0x30, 0x20)),
    ]
    for x1, y1, x2, y2, color in cross_edges:
        try:
            add_line_connector(slide, x1, y1, x2, y2, color, Pt(0.5))
        except:
            pass

    # Right panel
    add_rect(slide, Inches(9.6), Inches(0.9), Inches(3.4), Inches(5.8),
             RGBColor(0x12, 0x12, 0x18), RGBColor(0x30, 0x30, 0x40), Pt(1))

    add_multiline_text(slide, Inches(9.8), Inches(1.0), Inches(3.0), Inches(5.5), [
        ("Sigma.js v3 + Gephi", 14, ACCENT_YELLOW, True),
        ("ForceAtlas2 Layout", 11, TEXT_DIM),
        ("", 6, TEXT_DIM),
        ("Pipeline:", 11, TEXT_WHITE, True),
        ("  1. Export Neo4j → GEXF", 9, TEXT_DIM),
        ("  2. Gephi: ForceAtlas2", 9, TEXT_DIM),
        ("  3. Community detection", 9, TEXT_DIM),
        ("  4. Export → Sigma.js", 9, TEXT_DIM),
        ("", 6, TEXT_DIM),
        ("Strengths:", 11, TEXT_WHITE, True),
        ("  100K+ nodes (WebGL)", 9, ACCENT_GREEN),
        ("  Best for full-corpus view", 9, ACCENT_GREEN),
        ("  Academic SNA aesthetic", 9, ACCENT_GREEN),
        ("  Natural community layout", 9, ACCENT_GREEN),
        ("", 6, TEXT_DIM),
        ("Shown here:", 11, TEXT_WHITE, True),
        ("  7 communities = 7 Layers", 9, TEXT_DIM),
        ("  Dot size = degree centrality", 9, TEXT_DIM),
        ("  Cross-community edges =", 9, TEXT_DIM),
        ("    cross-layer conspirators", 9, TEXT_DIM),
        ("", 6, TEXT_DIM),
        ("Full corpus potential:", 11, ACCENT_YELLOW, True),
        ("  13,495 Record Nos as nodes", 9, TEXT_DIM),
        ("  Each citing atom = edge", 9, TEXT_DIM),
        ("  Galaxy-scale evidence map", 9, TEXT_DIM),
    ])


# ════════════════════════════════════════
# COMPARISON SLIDE — All 5 at a glance
# ════════════════════════════════════════
def slide_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.4), Inches(0.15), Inches(12), Inches(0.45),
                 "Visualization Comparison & Recommendation", 24, ACCENT_BLUE, True)
    add_rect(slide, Inches(0.4), Inches(0.55), Inches(12.5), Pt(1), ACCENT_BLUE)

    # Comparison table
    headers = ["", "Cytoscape.js", "Obsidian", "G6/AntV", "3D Force", "Sigma.js"]
    col_w = Inches(2.0)
    label_w = Inches(2.5)
    start_x = Inches(0.4)
    header_y = Inches(0.8)

    # Header row
    add_rect(slide, start_x, header_y, Inches(12.5), Inches(0.5), BG_CARD)
    colors_h = [TEXT_DIM, ACCENT_GREEN, ACCENT_PURPLE, ACCENT_ORANGE, ACCENT_CYAN, ACCENT_YELLOW]
    for i, (h, c) in enumerate(zip(headers, colors_h)):
        x = start_x + (Inches(0) if i == 0 else label_w + col_w * (i - 1))
        w = label_w if i == 0 else col_w
        add_text_box(slide, x, header_y + Inches(0.05), w, Inches(0.4),
                     h, 11, c, True, PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    rows = [
        ("Renderer", ["Canvas", "PIXI.js (WebGL)", "Canvas/WebGL", "Three.js (WebGL)", "WebGL"]),
        ("Dark Theme", ["Manual CSS", "Default dark", "Built-in toggle", "Default dark", "Config"]),
        ("Angular", ["Excellent", "N/A (replicate)", "Good (vanilla)", "Use vanilla lib", "Good (vanilla)"]),
        ("Max Nodes", ["~10K", "~10K", "~10K+", "~10K", "100K+"]),
        ("Rich Node Content", ["Basic labels", "Minimal dots", "Cards+badges", "Spheres", "Dots+labels"]),
        ("Role-Anchor Display", ["Border color", "Hover only", "Inline in card", "Hover tooltip", "On zoom"]),
        ("Visual Polish", ["Medium", "Iconic", "Highest", "Very High (3D)", "Academic"]),
        ("Best For", ["Production (v1)", "Aesthetic ref", "Production (v2)", "Presentations", "Full corpus"]),
    ]

    for ri, (label, values) in enumerate(rows):
        y = Inches(1.4) + ri * Inches(0.6)
        if ri % 2 == 0:
            add_rect(slide, start_x, y, Inches(12.5), Inches(0.5), BG_CARD)

        add_text_box(slide, start_x + Inches(0.1), y + Inches(0.05), label_w - Inches(0.2), Inches(0.4),
                     label, 10, TEXT_WHITE, True)
        for vi, val in enumerate(values):
            x = start_x + label_w + col_w * vi
            # Highlight recommended cells
            color = TEXT_WHITE
            if val in ["Built-in toggle", "Cards+badges", "Inline in card", "Highest", "Production (v2)"]:
                color = ACCENT_GREEN
            elif val in ["N/A (replicate)"]:
                color = TEXT_DIM
            add_text_box(slide, x, y + Inches(0.05), col_w, Inches(0.4),
                         val, 9, color, False, PP_ALIGN.CENTER)

    # Recommendation bar
    rec_y = Inches(6.3)
    add_rect(slide, Inches(0.4), rec_y, Inches(12.5), Inches(0.9), BG_PANEL, ACCENT_ORANGE, Pt(2))

    add_multiline_text(slide, Inches(0.7), rec_y + Inches(0.05), Inches(12), Inches(0.8), [
        ("Recommendation for Aurora v2:", 14, TEXT_BRIGHT, True),
        ("Primary: G6/AntV (production daily-use)  |  Secondary: Cytoscape.js (reuse v1 code, graph algorithms)  |  "
         "Presentation mode: 3D Force Graph  |  Full-corpus overview: Sigma.js/Gephi", 10, TEXT_DIM),
    ])


# ════════════════════════════════════════
# MAIN
# ════════════════════════════════════════
def main():
    out_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "output")
    pptx_path = os.path.join(out_dir, "aurora-v2-ui-mockup-2026-04-11.pptx")

    prs = Presentation(pptx_path)

    # Add new slides (after existing slide 9)
    slide_section_header(prs)    # Slide 10 — Section divider
    slide_variant_a_cytoscape(prs)  # Slide 11 — Cytoscape (v1 actual)
    slide_variant_b_obsidian(prs)   # Slide 12 — Obsidian constellation
    slide_variant_c_g6(prs)         # Slide 13 — G6/AntV
    slide_variant_d_3d(prs)         # Slide 14 — 3D Force Graph
    slide_variant_e_sigma(prs)      # Slide 15 — Sigma.js/Gephi
    slide_comparison(prs)           # Slide 16 — Comparison + recommendation

    prs.save(pptx_path)
    print(f"Updated: {pptx_path}")
    print(f"Total slides: {len(prs.slides)}")
    print("Added slides 10-16:")
    print("  10. Section header — 5 Visualization Approaches")
    print("  11. Variant A — Cytoscape.js (Aurora v1 actual style)")
    print("  12. Variant B — Obsidian Constellation Style")
    print("  13. Variant C — G6/AntV Enterprise Dark Theme [RECOMMENDED]")
    print("  14. Variant D — 3D Force Graph (Three.js) [WOW FACTOR]")
    print("  15. Variant E — Sigma.js / Gephi Pipeline [SCALE]")
    print("  16. Comparison & Recommendation")


if __name__ == "__main__":
    main()
